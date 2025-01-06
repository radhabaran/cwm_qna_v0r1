# document_processor.py

import os
from typing import List, Dict, Generator
import pypdf
from pptx import Presentation
from langchain.text_splitter import RecursiveCharacterTextSplitter
from openai import OpenAI
from tqdm import tqdm
from qdrant_client import QdrantClient
from qdrant_client.http import models


class DocumentProcessor:
    def __init__(self, config):
        """Initialize document processor with all necessary components"""
        self.config = config
        self.openai_client = OpenAI(api_key=config.OPENAI_API_KEY)
        
        # Initialize Qdrant client
        self.qdrant_client = QdrantClient(path=config.LOCAL_QDRANT_PATH)
        self._setup_collection()


        # Text splitter for chunking
        self.text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=config.CHUNK_SIZE,
            chunk_overlap=config.CHUNK_OVERLAP,
            length_function=len,
            separators=["\n\n", "\n", ".", " ", ""]
        )


    def _setup_collection(self):
        """Setup Qdrant collection if it doesn't exist"""
        collections = self.qdrant_client.get_collections()
        if not any(c.name == self.config.COLLECTION_NAME for c in collections.collections):
            self.qdrant_client.create_collection(
                collection_name=self.config.COLLECTION_NAME,
                vectors_config=models.VectorParams(
                    size=1536,
                    distance=models.Distance.COSINE
                )
            )


    def _get_processed_files(self) -> set:
        """Get set of already processed files"""
        try:
            response = self.qdrant_client.scroll(
                collection_name=self.config.COLLECTION_NAME,
                limit=10000,
                with_payload=['filename'],
                with_vectors=False
            )
            return {point.payload['filename'] for point in response[0]}
        except Exception:
            return set()


    def process_pdf(self, file_path: str) -> Generator[tuple, None, None]:
        """Process a single PDF file"""
        try:
            reader = pypdf.PdfReader(file_path)
            for page_number, page in enumerate(reader.pages):
                text = page.extract_text()
                if text.strip():  # Only yield if text is not empty
                    # Extract page header (first line)
                    lines = text.strip().split('\n', 1)
                    page_header = lines[0].strip() if len(lines) > 1 else ""
                    main_text = lines[1].strip() if len(lines) > 1 else lines[0].strip()
                    yield page_number, main_text, page_header

        except Exception as e:
            print(f"Error processing PDF {file_path}: {str(e)}")
            yield from []


    def process_pptx(self, file_path: str) -> Generator[tuple, None, None]:
        """Process a PPTX file"""
        try:
            presentation = Presentation(file_path)
            for slide_number, slide in enumerate(presentation.slides):
                slide_text = []
                slide_title = ""

                # Extract title if exists
                if slide.shapes.title:
                    slide_title = slide.shapes.title.text.strip()

                # Extract text from all shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text = shape.text.strip()
                        if text and text != slide_title:
                            slide_text.append(text)

                main_text = "\n".join(slide_text)
                if main_text or slide_title:
                    yield slide_number, main_text, slide_title

        except Exception as e:
            print(f"Error processing PPTX {file_path}: {str(e)}")
            yield from []


    def process_ppt(self, file_path: str) -> Generator[tuple, None, None]:
        """Process a legacy PPT file using python-pptx"""
        return self.process_pptx(file_path)  # Handle both .ppt and .pptx the same way


    def process_document(self, file_path: str) -> Generator[tuple, None, None]:
        """Process any supported document type"""
        file_extension = os.path.splitext(file_path)[1].lower()
        
        if file_extension == '.pdf':
            yield from self.process_pdf(file_path)
        elif file_extension == '.pptx':
            yield from self.process_pptx(file_path)
        elif file_extension == '.ppt':
            yield from self.process_ppt(file_path)
        else:
            print(f"Unsupported file type: {file_extension}")
            yield from []


    def create_chunks(self, text: str, metadata: Dict) -> List[Dict]:
        """Create chunks from text with metadata"""
        chunks = self.text_splitter.split_text(text)
        return [
            {
                'text': chunk,
                'metadata': {
                    'filename': metadata['filename'],
                    'page_number': metadata['page_number'],
                    'chunk_number': i + 1,
                    'page_header': metadata['page_header']
                }
            }
            for i, chunk in enumerate(chunks)
        ]


    def get_embeddings(self, texts: List[str]) -> List[List[float]]:
        """Get embeddings for texts using OpenAI"""
        try:
            response = self.openai_client.embeddings.create(
                model="text-embedding-ada-002",
                input=texts
            )
            return [data.embedding for data in response.data]
        except Exception as e:
            print(f"Error generating embeddings: {str(e)}")
            raise


    def store_vectors(self, vectors: List[List[float]], chunks: List[Dict]):
        """Store vectors and metadata in Qdrant"""
        try:
            points = [
                models.PointStruct(
                    id=abs(hash(f"{chunk['metadata']['filename']}_{chunk['metadata']['page_number']}_{chunk['metadata']['chunk_number']}")) % (2**63),
                    vector=vector,
                    payload={
                        'text': chunk['text'],
                        'filename': chunk['metadata']['filename'],
                        'page_number': chunk['metadata']['page_number'],
                        'chunk_number': chunk['metadata']['chunk_number'],
                        'page_header': chunk['metadata']['page_header']
                    }
                )
                for chunk, vector in zip(chunks, vectors)
            ]
            
            self.qdrant_client.upsert(
                collection_name=self.config.COLLECTION_NAME,
                points=points
            )
        except Exception as e:
            print(f"Error storing vectors: {str(e)}")
            raise


    def final_processing(self) -> int:
        """Main processing function"""
        # Get new files to process
        processed_files = self._get_processed_files()

        # Get files from both PDF and PPT directories
        pdf_files = {f for f in os.listdir(self.config.PDF_DIRECTORY) 
                    if f.endswith('.pdf')}
        ppt_files = {f for f in os.listdir(self.config.PPT_DIRECTORY) 
                    if f.endswith(('.ppt', '.pptx'))}
        
        current_files = pdf_files.union(ppt_files)
        new_files = current_files - processed_files

        if not new_files:
            print("No new files to process")
            return 0

        total_chunks = []
        total_processed = 0

        # Process each file
        for filename in tqdm(new_files, desc="Processing files"):
            if filename.endswith('.pdf'):
                file_path = os.path.join(self.config.PDF_DIRECTORY, filename)
            else:  # .ppt or .pptx
                file_path = os.path.join(self.config.PPT_DIRECTORY, filename)

            print(f"\nProcessing {filename}")
            
            # Process each page
            for page_number, text, page_header in self.process_document(file_path):
                if text:
                    chunks = self.create_chunks(text, {
                        'filename': filename,
                        'page_number': page_number + 1,
                        'page_header': page_header
                    })
                    total_chunks.extend(chunks)

                # Process in batches
                if len(total_chunks) >= self.config.BATCH_SIZE:
                    try:
                        batch = total_chunks[:self.config.BATCH_SIZE]
                        embeddings = self.get_embeddings([c['text'] for c in batch])
                        self.store_vectors(embeddings, batch)
                        total_processed += len(batch)
                        total_chunks = total_chunks[self.config.BATCH_SIZE:]
                        print(f"Processed {total_processed} chunks so far...")
                    except Exception as e:
                        print(f"Error processing batch: {str(e)}")
                        continue

        # Process remaining chunks
        if total_chunks:
            try:
                embeddings = self.get_embeddings([c['text'] for c in total_chunks])
                self.store_vectors(embeddings, total_chunks)
                total_processed += len(total_chunks)
            except Exception as e:
                print(f"Error processing final batch: {str(e)}")

        return total_processed


def main():
    from config import Config

    # Initialize directory structure
    directories = [
        Config.DOCUMENT_DIRECTORY,
        Config.PDF_DIRECTORY,
        Config.PPT_DIRECTORY,
        Config.LOCAL_QDRANT_PATH
    ]
    
    for directory in directories:
        os.makedirs(directory, exist_ok=True)
        
    processor = DocumentProcessor(Config)
    try:
        print("Starting document processing...")
        num_processed = processor.final_processing()
        print(f"Successfully processed {num_processed} chunks")
    except Exception as e:
        print(f"Error processing documents: {str(e)}")

if __name__ == "__main__":
    main()